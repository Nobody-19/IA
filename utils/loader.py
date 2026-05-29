"""
utils/loader.py
Chargement unifie de documents selon leur extension.
Remplace les deux versions dupliquees dans ingest.py et app.py.
"""

import re
import logging
from pathlib import Path

import fitz
from llama_index.core import Document

logger = logging.getLogger(__name__)

EXTENSIONS_SUPPORTEES = [
    ".pdf", ".docx", ".txt", ".xlsx", ".csv",
    ".pptx", ".md", ".rtf", ".odt"
]


def is_garbage_text(text: str) -> bool:
    if not text or len(text.strip()) < 20:
        return True
    weird = sum(
        1 for c in text
        if ord(c) > 127 and c not in "àâäéèêëîïôùûüçÀÂÄÉÈÊËÎÏÔÙÛÜÇ°€«»"
    )
    return (weird / len(text)) > 0.15


def clean_text(text: str) -> str:
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    text = re.sub(r" {3,}", "  ", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def _load_pdf(path: Path) -> Document | None:
    try:
        doc = fitz.open(str(path))
        text = ""
        for i, page in enumerate(doc):
            text += f"\n--- Page {i + 1} ---\n{page.get_text('text')}"
        doc.close()
        if is_garbage_text(text):
            logger.warning("PDF ignore (contenu illisible) : %s", path.name)
            return None
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": ".pdf"}
        )
    except Exception as e:
        logger.error("Erreur PDF %s : %s", path.name, e)
        return None


def _load_docx(path: Path) -> Document | None:
    try:
        from docx import Document as DocxDoc
        doc = DocxDoc(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
        if is_garbage_text(text):
            logger.warning("Word ignore : %s", path.name)
            return None
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": ".docx"}
        )
    except Exception as e:
        logger.error("Erreur DOCX %s : %s", path.name, e)
        return None


def _load_xlsx(path: Path) -> Document | None:
    try:
        import pandas as pd
        df = pd.read_excel(str(path))
        text = df.to_string()
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": ".xlsx"}
        )
    except Exception as e:
        logger.error("Erreur XLSX %s : %s", path.name, e)
        return None


def _load_csv(path: Path) -> Document | None:
    try:
        import pandas as pd
        df = pd.read_csv(str(path))
        text = df.to_string()
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": ".csv"}
        )
    except Exception as e:
        logger.error("Erreur CSV %s : %s", path.name, e)
        return None


def _load_pptx(path: Path) -> Document | None:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": ".pptx"}
        )
    except Exception as e:
        logger.error("Erreur PPTX %s : %s", path.name, e)
        return None


def _load_text(path: Path) -> Document | None:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
        if is_garbage_text(text):
            return None
        return Document(
            text=clean_text(text),
            metadata={"file_name": path.name, "file_path": str(path), "extension": path.suffix.lower()}
        )
    except Exception as e:
        logger.error("Erreur TXT %s : %s", path.name, e)
        return None


_LOADERS = {
    ".pdf":  _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".csv":  _load_csv,
    ".pptx": _load_pptx,
    ".txt":  _load_text,
    ".md":   _load_text,
    ".rtf":  _load_text,
    ".odt":  _load_text,
}


def load_doc(path: Path) -> Document | None:
    """
    Charge un document selon son extension.
    Retourne un Document LlamaIndex ou None si le fichier est illisible.
    """
    ext = path.suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        logger.warning("Extension non supportee : %s", path.suffix)
        return None
    return loader(path)


def charger_plusieurs(chemins: list[str]) -> list[Document]:
    """
    Charge une liste de chemins et retourne les Documents valides.
    """
    documents = []
    for chemin in chemins:
        doc = load_doc(Path(chemin))
        if doc:
            documents.append(doc)
        else:
            logger.warning("Document ignore : %s", chemin)
    return documents
